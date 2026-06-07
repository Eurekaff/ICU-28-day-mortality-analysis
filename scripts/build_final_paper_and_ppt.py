from __future__ import annotations

from pathlib import Path
from typing import Iterable

import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "outputs" / "docs"
TABLES = ROOT / "outputs" / "tables"
FIGURES = ROOT / "outputs" / "figures"

PAPER_OUT = DOCS / "基于MIMIC-IV早期结构化数据的ICU患者28天死亡风险预测研究_论文终稿.docx"
PPT_OUT = DOCS / "基于MIMIC-IV早期结构化数据的ICU患者28天死亡风险预测研究_汇报终稿.pptx"

TITLE = "基于MIMIC-IV早期结构化数据的ICU患者28天死亡风险预测研究"
SUBTITLE = "面向校准、缺失机制与解释稳定性的可复现实验"
EN_TITLE = (
    "Prediction of 28-Day Mortality in ICU Patients Using Early Structured MIMIC-IV Data: "
    "A Reproducible Study of Calibration, Missingness, Ablation, and Explanation Stability"
)


def fmt(x: float, digits: int = 3) -> str:
    return f"{float(x):.{digits}f}"


def pct(x: float, digits: int = 2) -> str:
    return f"{float(x) * 100:.{digits}f}%"


def read_csv(name: str) -> pd.DataFrame:
    return pd.read_csv(TABLES / name)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False, font_size: float = 8.0) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(str(text))
    r.bold = bold
    r.font.size = Pt(font_size)
    r.font.name = "Times New Roman"
    r._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc: Document, rows: list[list[str]], widths_cm: list[float] | None = None) -> None:
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    table.style = "Table Grid"
    tbl_pr = table._tbl.tblPr
    layout = OxmlElement("w:tblLayout")
    layout.set(qn("w:type"), "fixed")
    tbl_pr.append(layout)
    if widths_cm:
        tbl_w = OxmlElement("w:tblW")
        tbl_w.set(qn("w:type"), "dxa")
        tbl_w.set(qn("w:w"), str(sum(cm_to_dxa(w) for w in widths_cm)))
        tbl_pr.append(tbl_w)
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            set_cell_text(cell, value, bold=(i == 0), font_size=7.2)
            if i == 0:
                set_cell_shading(cell, "D9EAF7")
            if widths_cm:
                cell.width = Cm(widths_cm[j])
                tc_pr = cell._tc.get_or_add_tcPr()
                tc_w = tc_pr.first_child_found_in("w:tcW")
                if tc_w is None:
                    tc_w = OxmlElement("w:tcW")
                    tc_pr.append(tc_w)
                tc_w.set(qn("w:type"), "dxa")
                tc_w.set(qn("w:w"), str(cm_to_dxa(widths_cm[j])))
    doc.add_paragraph()


def cm_to_dxa(value: float) -> int:
    return int(value / 2.54 * 1440)


def set_font(run, size: float = 9, bold: bool = False, name: str = "宋体", color: str | None = None) -> None:
    run.font.size = Pt(size)
    run.bold = bold
    run.font.name = "Times New Roman"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    if color:
        run.font.color.rgb = RGBColor.from_string(color)


def add_paragraph(doc: Document, text: str, first_line: bool = True) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(3)
    p.paragraph_format.line_spacing = 1.08
    if first_line:
        p.paragraph_format.first_line_indent = Pt(18)
    r = p.add_run(text)
    set_font(r, 9, name="宋体")


def add_heading(doc: Document, text: str, level: int = 1) -> None:
    p = doc.add_paragraph()
    p.paragraph_format.keep_with_next = True
    p.paragraph_format.space_before = Pt(8 if level == 1 else 5)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    set_font(r, 10 if level == 1 else 9, bold=True, name="黑体")


def add_caption(doc: Document, zh: str, en: str, before: bool = True) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if before:
        p.paragraph_format.keep_with_next = True
    p.paragraph_format.space_after = Pt(3 if before else 6)
    r = p.add_run(f"{zh}\n{en}")
    set_font(r, 7, name="宋体")


def short_model(name: str) -> str:
    mapping = {
        "HistGradientBoosting": "HGB",
        "LogisticRegression": "Logistic",
        "RandomForest": "RF",
        "DecisionTree": "Tree",
    }
    return mapping.get(str(name), str(name))


def configure_base_doc(doc: Document) -> None:
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(1.8)
    section.bottom_margin = Cm(1.8)
    section.left_margin = Cm(1.7)
    section.right_margin = Cm(1.7)
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    styles["Normal"].font.size = Pt(9)


def set_two_columns(section) -> None:
    sect_pr = section._sectPr
    cols = sect_pr.xpath("./w:cols")
    cols = cols[0] if cols else OxmlElement("w:cols")
    cols.set(qn("w:num"), "2")
    cols.set(qn("w:space"), "480")
    if cols.getparent() is None:
        sect_pr.append(cols)


def add_title_block(doc: Document) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(TITLE)
    set_font(r, 16, bold=True, name="黑体")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("方明正")
    set_font(r, 9, bold=True, name="宋体")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("信息学院 软件工程，2022级")
    set_font(r, 8, name="仿宋")


def add_abstracts(doc: Document, summary: pd.Series, best: pd.Series, cal_best: pd.Series) -> None:
    abstract = (
        f"摘 要  ICU患者短期死亡风险高，已有基于MIMIC数据库的死亡预测研究多集中于模型性能比较，"
        f"但对概率校准、缺失机制、特征来源贡献和解释稳定性的系统分析相对不足。为提升研究的临床可用性和方法学规范性，"
        f"基于MIMIC-IV衍生结构化数据构建ICU患者28天死亡风险预测模型，并围绕校准、缺失模式和解释稳定性重构实验方案。"
        f"纳入ICU住院时长不少于24小时且结局明确的患者{int(summary.n_total)}例，其中28天内死亡{int(summary.n_events)}例，"
        f"事件率为{pct(summary.event_rate)}。预测变量包括人口学信息、首次ICU入科单元、入ICU后前24小时生命体征和实验室指标。"
        f"采用分层训练/测试划分和交叉验证比较Logistic回归、决策树、随机森林、ExtraTrees、HistGradientBoosting、XGBoost、"
        f"LightGBM和CatBoost等模型，并进一步开展特征组消融、缺失指示与缺失负担特征、概率校准和SHAP稳定性实验。"
        f"结果显示，基线模型中{best.model}取得最高测试集AUROC（{fmt(best.test_auroc)}），AUPRC为{fmt(best.test_auprc)}；"
        f"校准实验中{cal_best.model}经{cal_best.calibration}校准后Brier score最低（{fmt(cal_best.test_brier)}）。"
        f"消融结果表明实验室指标和24小时摘要统计是主要信息来源，缺失模式特征可进一步提升三类梯度提升模型的AUROC和AUPRC。"
        f"SHAP稳定性分析显示年龄、阴离子间隙、尿素氮、呼吸频率、血氧饱和度和ICU单元等特征在重复抽样中稳定出现。"
        f"研究表明，早期结构化ICU数据可用于28天死亡风险分层；相比单纯追求最高AUROC，综合校准、缺失机制、消融和解释稳定性的评估框架更有助于形成可复现、可解释且面向临床使用边界的预测模型。"
    )
    add_paragraph(doc, abstract, first_line=False)
    p = doc.add_paragraph()
    r = p.add_run("关键词：MIMIC-IV；ICU；28天死亡；机器学习；模型校准；缺失机制；SHAP")
    set_font(r, 9, bold=True, name="宋体")

    p = doc.add_paragraph()
    r = p.add_run(EN_TITLE)
    set_font(r, 12, bold=True, name="Times New Roman")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    english = (
        "Abstract  Short-term mortality prediction in the intensive care unit has been widely studied with public critical-care databases, "
        "yet many studies emphasize model discrimination while paying less attention to calibration, informative missingness, feature-source contribution, "
        "and explanation stability. This study developed a reproducible 28-day mortality prediction workflow using structured variables from the first 24 hours "
        "after ICU admission in a MIMIC-IV-derived cohort. Patients with an ICU length of stay of at least 24 hours and available 28-day mortality labels were included. "
        f"The final cohort contained {int(summary.n_total)} ICU records, including {int(summary.n_events)} deaths within 28 days. Candidate predictors included demographics, "
        "first care unit, vital signs, and laboratory measurements. Eight baseline models were compared using stratified train-test split and cross-validation, followed by "
        "additional experiments on feature ablation, missingness indicators, group-level missingness burden, probability calibration, risk stratification, and SHAP-based "
        f"explanation stability. {best.model} achieved the highest baseline test AUROC of {fmt(best.test_auroc)}, with an AUPRC of {fmt(best.test_auprc)}. "
        f"In the calibration analysis, {cal_best.model} with {cal_best.calibration} calibration achieved the lowest Brier score of {fmt(cal_best.test_brier)}. "
        "Ablation experiments indicated that laboratory variables and 24-hour summary features contributed substantial predictive information, while missingness-aware "
        "features improved discrimination and precision-recall performance across selected gradient boosting models. Stable SHAP features included age, anion gap, blood urea nitrogen, "
        "respiratory rate, oxygen saturation, and first care unit. These results suggest that early structured ICU data are informative for 28-day mortality prediction, and that "
        "clinically oriented evaluation should move beyond AUROC-only model selection toward calibration-aware, missingness-aware, and stability-aware assessment."
    )
    add_paragraph(doc, english, first_line=False)
    p = doc.add_paragraph()
    r = p.add_run("Keywords  MIMIC-IV; intensive care unit; 28-day mortality; machine learning; calibration; missingness; SHAP")
    set_font(r, 9, bold=True, name="Times New Roman")


def add_model_table(doc: Document, perf: pd.DataFrame) -> None:
    rows = [["模型", "AUROC", "AUPRC", "F1", "Brier"]]
    for _, row in perf.sort_values("test_auroc", ascending=False).iterrows():
        rows.append([
            short_model(row.model),
            fmt(row.test_auroc),
            fmt(row.test_auprc),
            fmt(row.test_f1),
            fmt(row.test_brier),
        ])
    add_caption(doc, "表1  候选模型在测试集上的预测性能", "Table 1  Predictive performance of candidate models on the test set")
    add_table(doc, rows, [2.5, 1.1, 1.1, 1.0, 1.0])


def add_ablation_table(doc: Document, ablation: pd.DataFrame) -> None:
    best = ablation.sort_values(["experiment", "test_auroc"], ascending=[True, False]).groupby("experiment").head(1)
    label = {
        "demographics": "人口学+ICU单元",
        "vitals": "生命体征",
        "labs": "实验室指标",
        "first_only": "首次测量值",
        "summary_only": "24小时摘要",
        "full": "完整特征",
    }
    order = ["demographics", "vitals", "labs", "first_only", "summary_only", "full"]
    rows = [["实验组", "模型", "AUROC", "AUPRC"]]
    for key in order:
        row = best[best.experiment == key].iloc[0]
        rows.append([label[key], short_model(row.model), fmt(row.test_auroc), fmt(row.test_auprc)])
    add_caption(doc, "表2  特征组消融实验结果", "Table 2  Feature-source ablation results")
    add_table(doc, rows, [2.2, 2.4, 1.1, 1.1])


def add_missing_table(doc: Document, miss: pd.DataFrame) -> None:
    rows = [["模型", "基础AUROC", "增强AUROC", "AUPRC增益"]]
    for model in ["HistGradientBoosting", "XGBoost", "LightGBM"]:
        base = miss[(miss.model == model) & (miss.experiment == "base_imputation")].iloc[0]
        aug = miss[(miss.model == model) & (miss.experiment == "missing_indicators_and_burden")].iloc[0]
        rows.append([
            short_model(model),
            fmt(base.test_auroc),
            fmt(aug.test_auroc),
            f"{aug.test_auprc - base.test_auprc:+.3f}",
        ])
    add_caption(doc, "表3  缺失机制增强实验结果", "Table 3  Results of missingness-aware feature augmentation")
    add_table(doc, rows, [2.4, 1.4, 1.4, 1.4])


def add_calibration_table(doc: Document, cal: pd.DataFrame) -> None:
    rows = [["模型", "校准", "AUROC", "Brier"]]
    for _, row in cal.sort_values("test_brier").head(6).iterrows():
        rows.append([short_model(row.model), row.calibration, fmt(row.test_auroc), fmt(row.test_brier)])
    add_caption(doc, "表4  概率校准实验结果（按Brier score排序）", "Table 4  Probability calibration results ranked by Brier score")
    add_table(doc, rows, [2.4, 1.4, 1.1, 1.1])


def add_references(doc: Document) -> None:
    refs = [
        "JOHNSON A E W, BULGARELLI L, SHEN L, et al. MIMIC-IV, a freely accessible electronic health record dataset[J]. Scientific Data, 2023, 10: 1.",
        "JOHNSON A E W, POLLARD T J, SHEN L, et al. MIMIC-III, a freely accessible critical care database[J]. Scientific Data, 2016, 3: 160035.",
        "HARUTYUNYAN H, KHACHATRIAN H, KALE D C, et al. Multitask learning and benchmarking with clinical time series data[J]. Scientific Data, 2019, 6: 96.",
        "KNAUS W A, DRAPER E A, WAGNER D P, et al. APACHE II: a severity of disease classification system[J]. Critical Care Medicine, 1985, 13(10): 818-829.",
        "LE GALL J R, LEMESHOW S, SAULNIER F. A new Simplified Acute Physiology Score (SAPS II) based on a European/North American multicenter study[J]. JAMA, 1993, 270(24): 2957-2963.",
        "VINCENT J L, MORENO R, TAKALA J, et al. The SOFA score to describe organ dysfunction/failure[J]. Intensive Care Medicine, 1996, 22(7): 707-710.",
        "COLLINS G S, REITSMA J B, ALTMAN D G, et al. Transparent Reporting of a multivariable prediction model for Individual Prognosis Or Diagnosis (TRIPOD)[J]. Annals of Internal Medicine, 2015, 162(1): 55-63.",
        "COLLINS G S, DHIMAN P, NAVARRO C L A, et al. TRIPOD+AI statement: updated guidance for reporting clinical prediction models that use regression or machine learning methods[J]. BMJ, 2024, 385: e078378.",
        "CHEN T, GUESTRIN C. XGBoost: a scalable tree boosting system[C]//Proceedings of the 22nd ACM SIGKDD International Conference on Knowledge Discovery and Data Mining. 2016: 785-794.",
        "KE G, MENG Q, FINLEY T, et al. LightGBM: a highly efficient gradient boosting decision tree[C]//Advances in Neural Information Processing Systems. 2017: 3146-3154.",
        "PROKHORENKOVA L, GUSEV G, VOROBEV A, et al. CatBoost: unbiased boosting with categorical features[C]//Advances in Neural Information Processing Systems. 2018: 6638-6648.",
        "LUNDBERG S M, LEE S I. A unified approach to interpreting model predictions[C]//Advances in Neural Information Processing Systems. 2017: 4765-4774.",
    ]
    add_heading(doc, "参考文献", 1)
    for i, ref in enumerate(refs, start=1):
        p = doc.add_paragraph()
        p.paragraph_format.first_line_indent = Pt(-18)
        p.paragraph_format.left_indent = Pt(18)
        r = p.add_run(f"[{i}] {ref}")
        set_font(r, 8, name="宋体")


def build_docx() -> None:
    summary = read_csv("icu_28d_mortality_dataset_summary.csv").iloc[0]
    perf = read_csv("icu_28d_mortality_table2_model_performance.csv")
    ablation = read_csv("icu_28d_mortality_ablation_performance.csv")
    miss = read_csv("icu_28d_mortality_missingness_performance.csv")
    cal = read_csv("icu_28d_mortality_calibration_comparison.csv")
    shap = read_csv("icu_28d_mortality_shap_stability.csv")
    bootstrap = read_csv("icu_28d_mortality_table4_best_model_bootstrap_ci.csv")
    best = perf.sort_values("test_auroc", ascending=False).iloc[0]
    cal_best = cal.sort_values("test_brier").iloc[0]

    doc = Document()
    configure_base_doc(doc)
    add_title_block(doc)
    add_abstracts(doc, summary, best, cal_best)

    body_section = doc.add_section(WD_SECTION.CONTINUOUS)
    set_two_columns(body_section)

    add_heading(doc, "1 引言", 1)
    add_paragraph(doc, "ICU患者病情变化快、短期死亡风险高，入科早期进行客观风险分层有助于识别高危患者、优化监护强度和支持资源配置。随着电子健康记录和开放重症数据库的发展，基于MIMIC数据的ICU预后预测研究逐渐增多。现有研究通常通过机器学习模型整合生命体征和实验室指标，并以AUROC作为主要评价指标。")
    add_paragraph(doc, "然而，类似研究已经较多，若仅停留在“使用MIMIC-IV、比较多个机器学习模型、用SHAP解释最佳模型”的路径上，创新性不足。临床风险模型不仅需要较好的区分能力，还需要概率校准可靠、缺失处理透明、特征贡献可解释、解释结果稳定。本研究据此将重点从单纯模型排行榜转向面向临床可用性的评估框架。")
    add_paragraph(doc, "研究贡献主要包括：第一，在基线模型比较基础上同步评价AUROC、AUPRC、Brier score、阈值指标和bootstrap不确定性；第二，将缺失指示和特征组缺失负担作为临床监测强度的代理信息纳入实验；第三，通过人口学、生命体征、实验室指标、首次值和24小时摘要特征消融回答模型信息来源；第四，采用重复抽样的SHAP稳定性分析，区分稳定预测特征和单次解释结果。")

    add_heading(doc, "2 相关工作", 1)
    add_paragraph(doc, "传统ICU评分系统如APACHE II、SAPS II和SOFA使用少量经过临床筛选的生理变量和固定权重描述病情严重程度，解释性强、应用历史长，但对电子病历中高维结构化变量和非线性关系的利用有限。")
    add_paragraph(doc, "MIMIC-IV为真实世界ICU电子病历研究提供了可复现数据基础。Harutyunyan等建立的临床时间序列benchmark强调任务定义、时间窗和可复现处理的重要性。与这些工作一致，本研究明确预测时点为ICU入科后24小时，结局为入科后28天死亡。")
    add_paragraph(doc, "TRIPOD和TRIPOD+AI强调预测模型研究需透明报告研究对象、预测时点、变量可得性、缺失处理、模型选择和验证方式。因此，本研究将创新点定位于规范化实验设计和可用性验证，而非声称提出全新的算法。")

    add_heading(doc, "3 资料与方法", 1)
    add_heading(doc, "3.1 数据来源与研究对象", 2)
    add_paragraph(doc, f"研究使用MIMIC-IV衍生分析数据集ch3_analysis_dataset。纳入标准为ICU住院时长可用、ICU住院时长不少于24小时且28天死亡结局可用。最终纳入{int(summary.n_total)}例ICU记录，其中28天内死亡{int(summary.n_events)}例，事件率为{pct(summary.event_rate)}。")
    add_heading(doc, "3.2 结局与预测变量", 2)
    add_paragraph(doc, "主要结局为ICU入科后28天内死亡，death_28d=1表示入科后28天内死亡，death_28d=0表示未在28天内死亡。预测变量包括年龄、性别、种族分组、首次ICU入科单元，以及入ICU后前24小时生命体征和实验室指标。数值特征包括first、mean、min和max等摘要统计，以表示早期状态和波动范围。")
    add_heading(doc, "3.3 建模流程", 2)
    add_paragraph(doc, "数据按8:2进行分层训练集/测试集划分。数值变量采用中位数填补并标准化，分类变量采用众数填补并one-hot编码。候选模型包括Logistic回归、决策树、随机森林、ExtraTrees、HistGradientBoosting、XGBoost、LightGBM和CatBoost。训练集内通过分层交叉验证选择模型参数，独立测试集用于最终评估。")
    add_heading(doc, "3.4 重构实验设计", 2)
    add_paragraph(doc, "为突出区别于常规模型比较研究的创新点，新增四类实验：（1）特征组消融，比较人口学、生命体征、实验室指标、首次测量值、24小时摘要和完整特征；（2）缺失机制实验，在基础填补上加入数值变量缺失指示器和特征组缺失负担；（3）概率校准实验，比较未校准、sigmoid校准和isotonic校准；（4）解释稳定性实验，对树模型进行重复抽样SHAP分析，统计Top 10特征出现频率。")

    add_heading(doc, "4 实验结果", 1)
    add_heading(doc, "4.1 基线模型性能", 2)
    add_paragraph(doc, f"如表1所列，{best.model}在基线模型中取得最高测试集AUROC（{fmt(best.test_auroc)}）和AUPRC（{fmt(best.test_auprc)}）。梯度提升类模型整体优于单棵决策树和传统线性模型，说明ICU早期结构化变量中存在明显的非线性预测信息。")
    add_model_table(doc, perf)
    add_paragraph(doc, f"Bootstrap分析显示，最佳模型AUROC均值为{fmt(bootstrap[bootstrap.metric == 'auroc'].iloc[0]['mean'])}，95%置信区间为{fmt(bootstrap[bootstrap.metric == 'auroc'].iloc[0]['lcl_95'])}-{fmt(bootstrap[bootstrap.metric == 'auroc'].iloc[0]['ucl_95'])}，提示模型区分能力较稳定。")

    add_heading(doc, "4.2 特征组消融", 2)
    add_paragraph(doc, "如表2所列，完整特征集取得最高AUROC；24小时摘要特征在特征数较少的情况下已接近完整特征表现，提示入ICU后24小时内的均值、极值等摘要统计具有较高信息密度。实验室指标优于生命体征和人口学特征，是主要预测信息来源之一。")
    add_ablation_table(doc, ablation)

    add_heading(doc, "4.3 缺失机制实验", 2)
    add_paragraph(doc, "如表3所列，引入缺失指示和特征组缺失负担后，HistGradientBoosting、XGBoost和LightGBM的AUROC与AUPRC均有提升。该结果说明ICU结构化数据中的缺失并非单纯噪声，检查是否发生以及不同特征组的缺失负担可能反映监测强度和临床关注点。")
    add_missing_table(doc, miss)

    add_heading(doc, "4.4 概率校准与风险分层", 2)
    add_paragraph(doc, f"如表4所列，按Brier score排序，{cal_best.model}经{cal_best.calibration}校准后概率可靠性最佳。LightGBM虽然在基线AUROC上最高，但未校准Brier score较高，经isotonic校准后Brier score明显下降，说明AUROC最高并不必然意味着概率最可靠。")
    add_calibration_table(doc, cal)

    add_heading(doc, "4.5 解释稳定性", 2)
    stable_features = ", ".join(shap["display_feature"].drop_duplicates().head(10).tolist())
    add_paragraph(doc, f"重复抽样SHAP稳定性分析显示，XGBoost和LightGBM的Top 10稳定特征高度一致，包括{stable_features}。这些变量分别对应基础脆弱性、代谢/肾功能、呼吸状态、氧合状态、ICU单元和部分数据记录模式。稳定性分析避免了将一次SHAP排序误读为稳定临床规律。")

    for fig_no, (path, zh, en) in enumerate([
        (FIGURES / "icu_28d_mortality_roc.png", "图1  不同模型ROC曲线", "Fig.1  ROC curves of candidate models"),
        (FIGURES / "icu_28d_mortality_pr.png", "图2  不同模型PR曲线", "Fig.2  Precision-recall curves of candidate models"),
        (FIGURES / "icu_28d_mortality_calibration.png", "图3  最佳模型校准曲线", "Fig.3  Calibration curve of the best baseline model"),
    ], start=1):
        if path.exists():
            doc.add_picture(str(path), width=Cm(7.2))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            add_caption(doc, zh, en, before=False)

    add_heading(doc, "5 讨论", 1)
    add_paragraph(doc, "本研究的核心发现不是某一算法本身具有首创性，而是通过更完整的实验框架揭示了ICU 28天死亡预测模型在区分能力、概率可靠性、数据处理机制和解释稳定性之间的差异。基线结果显示LightGBM区分度最高，但校准实验提示HistGradientBoosting和校准后的LightGBM在概率可靠性方面更具优势。")
    add_paragraph(doc, "缺失机制实验是本研究区别于常规中位数填补流程的重要部分。ICU场景中，某些检验项目是否被检测常由临床状态和治疗流程驱动，因此缺失模式可能携带风险信息。实验结果支持将缺失指示与缺失负担作为可解释的数据处理增强，但其含义应谨慎解释为监测强度或流程差异，而非病理因果关系。")
    add_paragraph(doc, "特征消融显示，实验室指标和24小时摘要统计贡献较大，符合短期预后受器官功能、代谢紊乱、氧合状态和感染炎症负荷共同影响的临床认知。解释稳定性分析进一步表明，年龄、BUN、阴离子间隙、呼吸频率和SpO2等特征在不同树模型和重复抽样中稳定出现，增强了模型解释的可信度。")
    add_paragraph(doc, "研究仍存在局限：第一，数据来自单一公开数据库衍生队列，尚缺少外部验证；第二，当前仅使用结构化变量，未纳入诊断、用药、治疗措施、护理记录和连续时间序列模型；第三，SHAP为模型归因工具，不能直接推断因果关系；第四，校准模型仍需在独立外部队列中验证后才能考虑临床部署。")

    add_heading(doc, "6 结束语", 1)
    add_paragraph(doc, "基于MIMIC-IV早期结构化数据的ICU 28天死亡预测模型具有较好的区分能力。通过引入特征消融、缺失机制增强、概率校准和SHAP稳定性分析，本研究将创新点从常规算法比较转向临床可用性和方法学规范性验证。后续工作应进一步开展外部验证、时间切分验证和前瞻性临床影响评估。")
    add_references(doc)

    doc.save(PAPER_OUT)


def add_ppt_textbox(slide, x, y, w, h, text, size=18, bold=False, color="1F2937", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.name = "PingFang SC"
    run.font.color.rgb = PptRGB.from_string(color)
    return box


def add_slide_title(slide, title: str, kicker: str | None = None) -> None:
    if kicker:
        add_ppt_textbox(slide, 0.55, 0.25, 5.0, 0.25, kicker, size=9, bold=True, color="2563EB")
    add_ppt_textbox(slide, 0.55, 0.48, 11.7, 0.45, title, size=22, bold=True, color="111827")
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.55), PptInches(1.05), PptInches(11.9), PptInches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGB.from_string("D1D5DB")
    shape.line.fill.background()


def add_bullets(slide, items: Iterable[str], x=0.75, y=1.35, w=5.6, h=4.8, size=15) -> None:
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(size)
        p.font.name = "PingFang SC"
        p.font.color.rgb = PptRGB.from_string("374151")
        p.space_after = PptPt(8)


def add_metric_card(slide, x, y, label, value, sub="") -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(2.55), PptInches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = PptRGB.from_string("EFF6FF")
    card.line.color.rgb = PptRGB.from_string("BFDBFE")
    add_ppt_textbox(slide, x + 0.15, y + 0.13, 2.25, 0.25, label, size=9, bold=True, color="1D4ED8")
    add_ppt_textbox(slide, x + 0.15, y + 0.42, 2.25, 0.38, value, size=21, bold=True, color="111827")
    if sub:
        add_ppt_textbox(slide, x + 0.15, y + 0.83, 2.25, 0.2, sub, size=8, color="4B5563")


def add_ppt_table(slide, rows: list[list[str]], x, y, w, h, font_size=9) -> None:
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    table = table_shape.table
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.margin_left = PptInches(0.04)
            cell.margin_right = PptInches(0.04)
            cell.margin_top = PptInches(0.03)
            cell.margin_bottom = PptInches(0.03)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = PptPt(font_size)
                    run.font.name = "PingFang SC"
                    run.font.bold = i == 0
                    run.font.color.rgb = PptRGB.from_string("111827")
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PptRGB.from_string("DBEAFE")


def build_pptx() -> None:
    summary = read_csv("icu_28d_mortality_dataset_summary.csv").iloc[0]
    perf = read_csv("icu_28d_mortality_table2_model_performance.csv")
    ablation = read_csv("icu_28d_mortality_ablation_performance.csv")
    miss = read_csv("icu_28d_mortality_missingness_performance.csv")
    cal = read_csv("icu_28d_mortality_calibration_comparison.csv")
    shap = read_csv("icu_28d_mortality_shap_stability.csv")
    best = perf.sort_values("test_auroc", ascending=False).iloc[0]
    cal_best = cal.sort_values("test_brier").iloc[0]

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptRGB.from_string("F8FAFC")
    add_ppt_textbox(slide, 0.8, 0.75, 10.8, 0.85, TITLE, size=28, bold=True, color="0F172A")
    add_ppt_textbox(slide, 0.82, 1.65, 9.8, 0.45, SUBTITLE, size=17, color="1D4ED8")
    add_ppt_textbox(slide, 0.85, 5.85, 4.5, 0.3, "方明正 | 信息学院 软件工程", size=12, color="475569")
    add_metric_card(slide, 0.85, 3.0, "队列规模", f"{int(summary.n_total):,}", "ICU记录")
    add_metric_card(slide, 3.65, 3.0, "死亡事件率", pct(summary.event_rate), f"{int(summary.n_events):,}例死亡")
    add_metric_card(slide, 6.45, 3.0, "最佳AUROC", fmt(best.test_auroc), best.model)
    add_metric_card(slide, 9.25, 3.0, "最佳Brier", fmt(cal_best.test_brier), f"{cal_best.model} + {cal_best.calibration}")

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "老师意见后的研究定位重构", "WHY")
    add_bullets(slide, [
        "类似研究很多：MIMIC-IV + 多模型比较 + SHAP 已经不是充分创新点",
        "本研究把重点从“哪个模型AUROC最高”转为“模型是否可校准、可解释、可复现”",
        "创新放在实验设计：缺失机制、特征消融、概率校准、解释稳定性",
        "答辩表述：不是声称新算法首创，而是强化临床预测模型的规范化评估",
    ], x=0.8, y=1.45, w=11.5, h=4.8, size=17)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "文献对照：已有工作与本研究差异", "LITERATURE")
    rows = [
        ["方向", "已有研究", "本研究回应"],
        ["ICU评分", "APACHE/SAPS/SOFA强调少量固定变量", "用早期结构化变量补充非线性信息"],
        ["MIMIC benchmark", "强调任务定义、时间窗和可复现", "明确24h预测时点与28天结局"],
        ["MIMIC死亡预测", "常见做法是GBDT + SHAP", "加入校准、缺失机制、消融和稳定性"],
        ["TRIPOD+AI", "强调透明报告和验证", "按规范重构实验与论文叙事"],
    ]
    add_ppt_table(slide, rows, 0.7, 1.35, 11.9, 4.65, font_size=12)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "整体实验设计", "METHOD")
    add_bullets(slide, [
        f"研究对象：ICU住院时长不少于24小时且28天结局明确，n={int(summary.n_total):,}",
        "预测时点：ICU入科后24小时；结局窗口：入科后28天死亡",
        "基线模型：Logistic、树模型、随机森林、ExtraTrees、HGB、XGBoost、LightGBM、CatBoost",
        "扩展实验：消融、缺失增强、校准、风险分层、SHAP稳定性",
    ], x=0.8, y=1.35, w=5.7, h=4.8, size=15)
    if (FIGURES / "icu_28d_mortality_roc.png").exists():
        slide.shapes.add_picture(str(FIGURES / "icu_28d_mortality_roc.png"), PptInches(6.8), PptInches(1.45), width=PptInches(5.55))

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "基线结果：LightGBM区分度最高", "RESULT 1")
    top = perf.sort_values("test_auroc", ascending=False).head(5)
    rows = [["模型", "AUROC", "AUPRC", "F1", "Brier"]]
    for _, r in top.iterrows():
        rows.append([r.model, fmt(r.test_auroc), fmt(r.test_auprc), fmt(r.test_f1), fmt(r.test_brier)])
    add_ppt_table(slide, rows, 0.7, 1.35, 6.0, 3.0, font_size=12)
    add_bullets(slide, [
        f"LightGBM AUROC={fmt(best.test_auroc)}，AUPRC={fmt(best.test_auprc)}",
        "HGB/XGBoost与LightGBM差距很小",
        "Brier显示概率可靠性与AUROC排序并不完全一致",
    ], x=7.2, y=1.45, w=5.3, h=3.5, size=15)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "创新实验一：特征组消融回答“性能来自哪里”", "ABLATION")
    best_ab = ablation.sort_values(["experiment", "test_auroc"], ascending=[True, False]).groupby("experiment").head(1)
    label = {"demographics": "人口学", "vitals": "生命体征", "labs": "实验室", "first_only": "首次值", "summary_only": "24h摘要", "full": "完整"}
    rows = [["实验组", "模型", "AUROC", "AUPRC"]]
    for key in ["demographics", "vitals", "labs", "first_only", "summary_only", "full"]:
        r = best_ab[best_ab.experiment == key].iloc[0]
        rows.append([label[key], r.model, fmt(r.test_auroc), fmt(r.test_auprc)])
    add_ppt_table(slide, rows, 0.75, 1.25, 7.0, 4.7, font_size=12)
    add_bullets(slide, [
        "完整特征最优，但24h摘要已接近完整特征",
        "实验室指标强于生命体征和人口学",
        "叙事价值：不只给模型结果，也解释信息来源",
    ], x=8.1, y=1.55, w=4.5, h=3.6, size=15)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "创新实验二：缺失模式被作为信息，而非只当作噪声", "MISSINGNESS")
    rows = [["模型", "基础AUROC", "缺失增强AUROC", "AUPRC增益"]]
    for model in ["HistGradientBoosting", "XGBoost", "LightGBM"]:
        base = miss[(miss.model == model) & (miss.experiment == "base_imputation")].iloc[0]
        aug = miss[(miss.model == model) & (miss.experiment == "missing_indicators_and_burden")].iloc[0]
        rows.append([model, fmt(base.test_auroc), fmt(aug.test_auroc), f"{aug.test_auprc - base.test_auprc:+.3f}"])
    add_ppt_table(slide, rows, 0.8, 1.35, 7.0, 3.1, font_size=12)
    add_bullets(slide, [
        "三类梯度提升模型AUROC均有提升",
        "AUPRC提升更明显，尤其HGB +0.014",
        "临床含义：检查是否发生可反映监测强度和临床关注点",
    ], x=8.2, y=1.45, w=4.5, h=3.8, size=15)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "创新实验三：校准后再看概率可靠性", "CALIBRATION")
    rows = [["模型", "校准", "AUROC", "Brier", "F1"]]
    for _, r in cal.sort_values("test_brier").head(6).iterrows():
        rows.append([r.model, r.calibration, fmt(r.test_auroc), fmt(r.test_brier), fmt(r.test_f1)])
    add_ppt_table(slide, rows, 0.75, 1.25, 7.4, 4.4, font_size=11)
    add_bullets(slide, [
        f"最佳概率可靠性：{cal_best.model} + {cal_best.calibration}",
        "LightGBM未校准Brier偏高，isotonic后明显改善",
        "答辩重点：AUROC最高并不等于风险概率最可靠",
    ], x=8.45, y=1.45, w=4.2, h=3.8, size=15)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "创新实验四：解释稳定性，而不是单次SHAP排名", "SHAP")
    stable = shap["display_feature"].drop_duplicates().head(10).tolist()
    add_bullets(slide, stable, x=0.9, y=1.25, w=4.6, h=5.3, size=14)
    if (FIGURES / "icu_28d_mortality_shap_importance_bar.png").exists():
        slide.shapes.add_picture(str(FIGURES / "icu_28d_mortality_shap_importance_bar.png"), PptInches(5.7), PptInches(1.35), width=PptInches(6.6))

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "论文创新点表达", "CONTRIBUTION")
    add_bullets(slide, [
        "评价框架创新：AUROC + AUPRC + Brier + 校准 + 阈值指标",
        "数据处理创新：缺失指示器 + 特征组缺失负担",
        "实验设计创新：人口学/生命体征/实验室/时间摘要消融",
        "解释验证创新：跨模型、重复抽样的SHAP稳定性",
        "报告规范创新：按TRIPOD+AI明确预测时点、变量可得性和局限",
    ], x=0.85, y=1.35, w=11.5, h=5.2, size=18)

    slide = prs.slides.add_slide(blank)
    add_slide_title(slide, "结论与下一步", "CONCLUSION")
    add_bullets(slide, [
        "早期结构化ICU数据可较好预测28天死亡风险",
        "创新点已经能从实验结果中得到支撑，不需要硬说“新模型首创”",
        "正式论文应突出：可复现流程、缺失机制、校准可靠性、解释稳定性",
        "局限：仍需外部验证、时间切分验证和前瞻性临床影响评估",
    ], x=0.9, y=1.35, w=11.3, h=4.8, size=18)

    prs.save(PPT_OUT)


def main() -> None:
    DOCS.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()
    print(PAPER_OUT)
    print(PPT_OUT)


if __name__ == "__main__":
    main()
